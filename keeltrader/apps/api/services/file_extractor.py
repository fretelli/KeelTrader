"""File content extraction service for various file types."""

import csv
import io
import json
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Optional, Tuple

from core.logging import get_logger

logger = get_logger(__name__)


@dataclass
class ExtractionResult:
    """Result of file content extraction."""

    success: bool
    text: Optional[str] = None
    error: Optional[str] = None
    file_type: Optional[str] = None
    page_count: Optional[int] = None


# File type categories
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".ico"}
AUDIO_EXTENSIONS = {".wav", ".mp3", ".ogg", ".webm", ".m4a", ".aac", ".flac"}
PDF_EXTENSIONS = {".pdf"}
WORD_EXTENSIONS = {".docx", ".doc"}
EXCEL_EXTENSIONS = {".xlsx", ".xls", ".csv"}
PPT_EXTENSIONS = {".pptx", ".ppt"}
TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".rst", ".log", ".ini", ".cfg", ".conf"}
JSON_EXTENSIONS = {".json", ".jsonl"}
XML_EXTENSIONS = {".xml", ".html", ".htm", ".xhtml", ".svg"}
YAML_EXTENSIONS = {".yaml", ".yml"}
CODE_EXTENSIONS = {
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".vue",
    ".svelte",
    ".java",
    ".c",
    ".cpp",
    ".cc",
    ".h",
    ".hpp",
    ".cs",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".swift",
    ".kt",
    ".scala",
    ".sh",
    ".bash",
    ".zsh",
    ".fish",
    ".ps1",
    ".bat",
    ".cmd",
    ".sql",
    ".r",
    ".m",
    ".pl",
    ".lua",
    ".dart",
    ".elm",
    ".css",
    ".scss",
    ".sass",
    ".less",
    ".styl",
}

# Binary files that cannot be extracted
BINARY_EXTENSIONS = {
    ".exe",
    ".dll",
    ".so",
    ".dylib",
    ".bin",
    ".zip",
    ".tar",
    ".gz",
    ".rar",
    ".7z",
    ".bz2",
    ".mp4",
    ".avi",
    ".mov",
    ".mkv",
    ".wmv",
    ".flv",
    ".iso",
    ".dmg",
    ".msi",
}


def get_file_category(filename: str) -> str:
    """
    Determine the category of a file based on its extension.

    Args:
        filename: Name of the file

    Returns:
        Category string: 'image', 'audio', 'pdf', 'word', 'excel', 'ppt',
                        'text', 'code', 'binary', or 'unknown'
    """
    ext = Path(filename).suffix.lower()

    if ext in IMAGE_EXTENSIONS:
        return "image"
    elif ext in AUDIO_EXTENSIONS:
        return "audio"
    elif ext in PDF_EXTENSIONS:
        return "pdf"
    elif ext in WORD_EXTENSIONS:
        return "word"
    elif ext in EXCEL_EXTENSIONS:
        return "excel"
    elif ext in PPT_EXTENSIONS:
        return "ppt"
    elif (
        ext in TEXT_EXTENSIONS
        or ext in JSON_EXTENSIONS
        or ext in XML_EXTENSIONS
        or ext in YAML_EXTENSIONS
    ):
        return "text"
    elif ext in CODE_EXTENSIONS:
        return "code"
    elif ext in BINARY_EXTENSIONS:
        return "binary"
    else:
        return "unknown"


def can_extract_text(filename: str) -> bool:
    """
    Check if text can be extracted from a file.

    Args:
        filename: Name of the file

    Returns:
        True if text extraction is possible
    """
    category = get_file_category(filename)
    return category in ("pdf", "word", "excel", "ppt", "text", "code")


async def extract_text(file_path: Path, filename: str) -> ExtractionResult:
    """
    Extract text content from a file.

    Args:
        file_path: Path to the file
        filename: Original filename (used for extension detection)

    Returns:
        ExtractionResult with extracted text or error
    """
    category = get_file_category(filename)

    try:
        if category == "pdf":
            return await _extract_pdf(file_path)
        elif category == "word":
            return await _extract_word(file_path)
        elif category == "excel":
            return await _extract_excel(file_path, filename)
        elif category == "ppt":
            return await _extract_ppt(file_path)
        elif category in ("text", "code"):
            return await _extract_text_file(file_path)
        else:
            return ExtractionResult(
                success=False,
                error=f"Cannot extract text from {category} files",
                file_type=category,
            )
    except Exception as e:
        logger.error(f"Failed to extract text from {filename}: {e}")
        return ExtractionResult(success=False, error=str(e), file_type=category)


async def _extract_pdf(file_path: Path) -> ExtractionResult:
    """Extract text from PDF file."""
    try:
        import PyPDF2
    except ImportError:
        try:
            import pypdf as PyPDF2
        except ImportError:
            return ExtractionResult(
                success=False,
                error="PyPDF2 library not installed. Run: pip install PyPDF2",
                file_type="pdf",
            )

    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            page_count = len(reader.pages)
            text_parts = []

            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- Page {i + 1} ---\n{page_text}")

            text = "\n\n".join(text_parts)

            return ExtractionResult(
                success=True, text=text, file_type="pdf", page_count=page_count
            )
    except Exception as e:
        return ExtractionResult(
            success=False, error=f"Failed to parse PDF: {str(e)}", file_type="pdf"
        )


async def _extract_word(file_path: Path) -> ExtractionResult:
    """Extract text from Word document."""
    try:
        from docx import Document
    except ImportError:
        return ExtractionResult(
            success=False,
            error="python-docx library not installed. Run: pip install python-docx",
            file_type="word",
        )

    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)

        # Also extract tables
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_rows.append(" | ".join(cells))
            if table_rows:
                text += "\n\n[Table]\n" + "\n".join(table_rows)

        return ExtractionResult(success=True, text=text, file_type="word")
    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"Failed to parse Word document: {str(e)}",
            file_type="word",
        )


async def _extract_excel(file_path: Path, filename: str) -> ExtractionResult:
    """Extract text from Excel file or CSV."""
    ext = Path(filename).suffix.lower()

    if ext == ".csv":
        return await _extract_csv(file_path)

    try:
        from openpyxl import load_workbook
    except ImportError:
        return ExtractionResult(
            success=False,
            error="openpyxl library not installed. Run: pip install openpyxl",
            file_type="excel",
        )

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"=== Sheet: {sheet_name} ===")

            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):  # Skip empty rows
                    rows.append(" | ".join(cells))

            text_parts.append("\n".join(rows))

        wb.close()
        text = "\n\n".join(text_parts)

        return ExtractionResult(success=True, text=text, file_type="excel")
    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"Failed to parse Excel file: {str(e)}",
            file_type="excel",
        )


async def _extract_csv(file_path: Path) -> ExtractionResult:
    """Extract text from CSV file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = []
            for row in reader:
                rows.append(" | ".join(row))
            text = "\n".join(rows)

        return ExtractionResult(success=True, text=text, file_type="csv")
    except Exception as e:
        return ExtractionResult(
            success=False, error=f"Failed to parse CSV file: {str(e)}", file_type="csv"
        )


async def _extract_ppt(file_path: Path) -> ExtractionResult:
    """Extract text from PowerPoint presentation."""
    try:
        from pptx import Presentation
    except ImportError:
        return ExtractionResult(
            success=False,
            error="python-pptx library not installed. Run: pip install python-pptx",
            file_type="ppt",
        )

    try:
        prs = Presentation(file_path)
        text_parts = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- Slide {i + 1} ---")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            text_parts.append("\n".join(slide_text))

        text = "\n\n".join(text_parts)

        return ExtractionResult(
            success=True, text=text, file_type="ppt", page_count=len(prs.slides)
        )
    except Exception as e:
        return ExtractionResult(
            success=False,
            error=f"Failed to parse PowerPoint: {str(e)}",
            file_type="ppt",
        )


async def _extract_text_file(file_path: Path) -> ExtractionResult:
    """Extract text from plain text or code file."""
    try:
        # Try UTF-8 first, then fall back to latin-1
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        except UnicodeDecodeError:
            with open(file_path, "r", encoding="latin-1") as f:
                text = f.read()

        return ExtractionResult(success=True, text=text, file_type="text")
    except Exception as e:
        return ExtractionResult(
            success=False, error=f"Failed to read text file: {str(e)}", file_type="text"
        )


def get_file_size_limit(filename: str) -> int:
    """
    Get the size limit for a file based on its type.

    Args:
        filename: Name of the file

    Returns:
        Maximum file size in bytes
    """
    category = get_file_category(filename)

    limits = {
        "image": 10 * 1024 * 1024,  # 10MB
        "audio": 25 * 1024 * 1024,  # 25MB
        "pdf": 50 * 1024 * 1024,  # 50MB
        "word": 50 * 1024 * 1024,  # 50MB
        "excel": 50 * 1024 * 1024,  # 50MB
        "ppt": 50 * 1024 * 1024,  # 50MB
        "text": 10 * 1024 * 1024,  # 10MB
        "code": 10 * 1024 * 1024,  # 10MB
        "binary": 100 * 1024 * 1024,  # 100MB
        "unknown": 100 * 1024 * 1024,  # 100MB
    }

    return limits.get(category, 100 * 1024 * 1024)


def get_allowed_mime_types() -> dict:
    """
    Get a mapping of file categories to allowed MIME types.

    Returns:
        Dictionary mapping categories to lists of MIME types
    """
    return {
        "image": [
            "image/jpeg",
            "image/png",
            "image/gif",
            "image/webp",
            "image/bmp",
            "image/x-icon",
        ],
        "audio": [
            "audio/wav",
            "audio/x-wav",
            "audio/mp3",
            "audio/mpeg",
            "audio/ogg",
            "audio/webm",
            "audio/m4a",
            "audio/aac",
            "audio/flac",
            "audio/x-flac",
        ],
        "pdf": ["application/pdf"],
        "word": [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ],
        "excel": [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
            "text/csv",
        ],
        "ppt": [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ],
        "text": [
            "text/plain",
            "text/markdown",
            "text/x-markdown",
            "text/html",
            "text/xml",
            "application/xml",
            "application/json",
            "text/yaml",
            "application/x-yaml",
        ],
        "code": [
            "text/x-python",
            "application/javascript",
            "text/javascript",
            "text/typescript",
            "text/x-java-source",
            "text/x-c",
            "text/x-c++",
            "text/x-csharp",
            "text/x-go",
            "text/x-rust",
            "text/x-ruby",
            "application/x-php",
            "text/x-swift",
            "text/x-kotlin",
            "text/x-scala",
            "text/x-shellscript",
            "text/css",
            "text/x-scss",
            "text/x-sass",
        ],
    }
